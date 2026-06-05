from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
GOLD_DARK   = RGBColor(0xB8, 0x86, 0x0B)   # #B8860B  dark-gold
GOLD_MID    = RGBColor(0xD4, 0xAF, 0x37)   # #D4AF37  classic gold
GOLD_LIGHT  = RGBColor(0xF5, 0xE6, 0xA3)   # #F5E6A3  pale gold accent
BG_CREAM    = RGBColor(0xFF, 0xFC, 0xF0)   # #FFFCF0  warm ivory
NAVY        = RGBColor(0x1A, 0x23, 0x4E)   # #1A234E  deep navy
DARK_GRAY   = RGBColor(0x2E, 0x2E, 0x2E)   # #2E2E2E  body text
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_TEAL = RGBColor(0x00, 0x6B, 0x6B)   # subtle teal for variety

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(2)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name="Palatino Linotype", font_size=Pt(18),
                bold=False, italic=False, color=DARK_GRAY,
                align=PP_ALIGN.LEFT, wrap=True, line_spacing=None):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name      = font_name
    run.font.size      = font_size
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    if line_spacing:
        from pptx.util import Pt as Pt2
        from pptx.oxml.ns import qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing.pt * 100)))
    return txb


def gold_border_frame(slide, thickness=Pt(3)):
    """Outer gold border rectangle around the full slide."""
    m = Inches(0.12)
    add_rect(slide,
             m, m,
             SLIDE_W - 2*m, SLIDE_H - 2*m,
             fill_color=None,
             line_color=GOLD_MID,
             line_width=thickness)


def inner_gold_border(slide, thickness=Pt(1.2)):
    """Thin inner decorative border."""
    m = Inches(0.22)
    add_rect(slide,
             m, m,
             SLIDE_W - 2*m, SLIDE_H - 2*m,
             fill_color=None,
             line_color=GOLD_LIGHT,
             line_width=thickness)


def corner_ornament(slide):
    """Small gold squares at corners for a royal look."""
    sz = Inches(0.15)
    positions = [
        (Inches(0.07), Inches(0.07)),
        (SLIDE_W - Inches(0.07) - sz, Inches(0.07)),
        (Inches(0.07), SLIDE_H - Inches(0.07) - sz),
        (SLIDE_W - Inches(0.07) - sz, SLIDE_H - Inches(0.07) - sz),
    ]
    for (l, t) in positions:
        add_rect(slide, l, t, sz, sz,
                 fill_color=GOLD_DARK,
                 line_color=None,
                 line_width=Pt(0))


def header_band(slide, subtitle=""):
    """Gold gradient-like header band (solid gold top stripe)."""
    band_h = Inches(1.45)
    add_rect(slide,
             Inches(0.0), Inches(0.0),
             SLIDE_W, band_h,
             fill_color=NAVY,
             line_color=None,
             line_width=Pt(0))
    # gold accent line under header
    add_rect(slide,
             Inches(0.0), band_h,
             SLIDE_W, Pt(4),
             fill_color=GOLD_MID,
             line_color=None,
             line_width=Pt(0))
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.45), Inches(1.1),
                    Inches(12.5), Inches(0.4),
                    font_size=Pt(11), color=GOLD_LIGHT,
                    italic=True, align=PP_ALIGN.LEFT)


def slide_title(slide, title_text):
    add_textbox(slide, title_text,
                Inches(0.45), Inches(0.15),
                Inches(12.4), Inches(0.9),
                font_name="Palatino Linotype",
                font_size=Pt(26), bold=True,
                color=GOLD_MID, align=PP_ALIGN.LEFT)


def footer_band(slide, slide_num, total=22):
    foot_h = Inches(0.35)
    top    = SLIDE_H - foot_h
    add_rect(slide,
             Inches(0.0), top,
             SLIDE_W, foot_h,
             fill_color=NAVY,
             line_color=None,
             line_width=Pt(0))
    add_textbox(slide,
                "Teacher Training  •  From Effort to Excellence",
                Inches(0.45), top + Pt(4),
                Inches(10), Inches(0.3),
                font_size=Pt(9), color=GOLD_LIGHT,
                italic=True, align=PP_ALIGN.LEFT)
    add_textbox(slide,
                f"{slide_num} / {total}",
                Inches(11.8), top + Pt(4),
                Inches(1.3), Inches(0.3),
                font_size=Pt(9), color=GOLD_LIGHT,
                align=PP_ALIGN.RIGHT)


def decorate(slide, slide_num, subtitle="", total=22):
    set_bg(slide, BG_CREAM)
    gold_border_frame(slide)
    inner_gold_border(slide)
    corner_ornament(slide)
    header_band(slide, subtitle)
    footer_band(slide, slide_num, total)


def content_top():
    return Inches(1.65)


def add_bullet_block(slide, items, left, top, width, height,
                     font_size=Pt(15.5), indent=False, color=DARK_GRAY,
                     bullet_char="✦", gap=Pt(28)):
    from pptx.oxml.ns import qn
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        prefix = f"    {bullet_char}  " if indent else f"{bullet_char}  "
        run.text = prefix + item
        run.font.name      = "Palatino Linotype"
        run.font.size      = font_size
        run.font.color.rgb = color
        # paragraph spacing
        pPr = p._p.get_or_add_pPr()
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', str(int(gap.pt * 100)))
    return txb


def gold_divider(slide, top, width_frac=0.85):
    w = SLIDE_W * width_frac
    l = (SLIDE_W - w) / 2
    add_rect(slide, l, top, w, Pt(2),
             fill_color=GOLD_MID,
             line_color=None,
             line_width=Pt(0))


def highlight_box(slide, text, left, top, width, height,
                  bg=NAVY, fg=GOLD_LIGHT, font_size=Pt(14), bold=False):
    add_rect(slide, left, top, width, height,
             fill_color=bg, line_color=GOLD_MID, line_width=Pt(1.5))
    add_textbox(slide, text,
                left + Inches(0.1), top + Inches(0.08),
                width - Inches(0.2), height - Inches(0.1),
                font_size=font_size, color=fg,
                bold=bold, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 1  –  Cover
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, NAVY)

# Outer gold border
m = Inches(0.12)
add_rect(slide, m, m, SLIDE_W-2*m, SLIDE_H-2*m,
         fill_color=None, line_color=GOLD_MID, line_width=Pt(4))
add_rect(slide, Inches(0.28), Inches(0.28),
         SLIDE_W-Inches(0.56), SLIDE_H-Inches(0.56),
         fill_color=None, line_color=GOLD_LIGHT, line_width=Pt(1.2))
corner_ornament(slide)

# Gold horizontal bands
for y_frac in [0.30, 0.72]:
    add_rect(slide, Inches(0), SLIDE_H*y_frac,
             SLIDE_W, Pt(3),
             fill_color=GOLD_MID, line_color=None, line_width=Pt(0))

# School / org name
add_textbox(slide, "TEACHER TRAINING PROGRAMME",
            Inches(1), Inches(0.55), Inches(11.33), Inches(0.55),
            font_name="Palatino Linotype", font_size=Pt(16),
            bold=False, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

# Main Title
add_textbox(slide, "From Effort to Excellence",
            Inches(0.8), Inches(1.5), Inches(11.73), Inches(1.2),
            font_name="Palatino Linotype", font_size=Pt(46),
            bold=True, color=GOLD_MID, align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, "A Roadmap to Academic Excellence and Outstanding Results",
            Inches(1.2), Inches(2.75), Inches(10.93), Inches(0.75),
            font_name="Palatino Linotype", font_size=Pt(19),
            italic=True, color=WHITE, align=PP_ALIGN.CENTER)

gold_divider(slide, Inches(3.65), 0.7)

# Key quote
add_textbox(slide, '"Excellence is never an accident; it is the result of consistent effort and strategic action."',
            Inches(1.5), Inches(3.85), Inches(10.33), Inches(0.8),
            font_name="Palatino Linotype", font_size=Pt(13.5),
            italic=True, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

# Three pillars
for i, (icon, label) in enumerate([("📚", "Accountability"), ("🎯", "Execution"), ("📊", "Outcomes")]):
    lx = Inches(2.5 + i * 3.1)
    add_rect(slide, lx, Inches(4.9), Inches(2.5), Inches(1.0),
             fill_color=RGBColor(0x26, 0x33, 0x66),
             line_color=GOLD_MID, line_width=Pt(1.5))
    add_textbox(slide, f"{icon}  {label}",
                lx + Inches(0.1), Inches(4.95),
                Inches(2.3), Inches(0.85),
                font_size=Pt(15.5), bold=True, color=GOLD_MID,
                align=PP_ALIGN.CENTER)

# Footer note
add_textbox(slide, "Building a Result-Oriented Teaching Culture  •  All Streams",
            Inches(1), Inches(6.95), Inches(11.33), Inches(0.4),
            font_size=Pt(10), color=GOLD_LIGHT,
            italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 2  –  Why Results Matter
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 2, "Foundation: Why We Focus on Results")
slide_title(slide, "Why Results Matter")

bullets = [
    "Results reflect the true effectiveness of teaching and learning.",
    "Academic excellence builds institutional credibility and trust.",
    "Every child can improve with the right, timely intervention.",
    "Results are not produced in March — they are built throughout the year.",
]
add_bullet_block(slide, bullets,
                 Inches(0.55), content_top(),
                 Inches(12.2), Inches(3.0),
                 font_size=Pt(16))

gold_divider(slide, Inches(4.85))

highlight_box(slide,
              '"Excellence is never an accident; it is the result of consistent effort and strategic action."',
              Inches(0.8), Inches(5.05),
              Inches(11.73), Inches(0.95),
              font_size=Pt(14.5), bold=False)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 3  –  Formula of Success
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 3, "The Core Formula")
slide_title(slide, "Understanding the Formula of Success")

components = [
    ("Effective Teaching",    "🎓"),
    ("Consistent Assessment", "📋"),
    ("Timely Intervention",   "⏱"),
    ("Parent Partnership",    "🤝"),
    ("Student Accountability","🏆"),
]

box_w = Inches(2.3)
box_h = Inches(1.55)
gap   = Inches(0.18)
total_w = len(components)*box_w + (len(components)-1)*gap
start_l = (SLIDE_W - total_w) / 2

for i, (label, icon) in enumerate(components):
    lx = start_l + i*(box_w + gap)
    add_rect(slide, lx, content_top(),
             box_w, box_h,
             fill_color=RGBColor(0xF0, 0xE8, 0xC8),
             line_color=GOLD_DARK, line_width=Pt(2))
    add_textbox(slide, icon,
                lx, content_top() + Inches(0.1),
                box_w, Inches(0.55),
                font_size=Pt(22), align=PP_ALIGN.CENTER)
    add_textbox(slide, label,
                lx, content_top() + Inches(0.65),
                box_w, Inches(0.8),
                font_size=Pt(12.5), bold=True, color=NAVY,
                align=PP_ALIGN.CENTER)

gold_divider(slide, content_top() + box_h + Inches(0.2))

# Visual formula
formula = "Achievement  =  Teaching Quality  ×  Student Effort  ×  Monitoring  ×  Support"
highlight_box(slide, formula,
              Inches(0.6), Inches(3.55),
              Inches(12.13), Inches(0.75),
              font_size=Pt(15), bold=True)

add_textbox(slide, "Student Success = Effective Teaching + Consistent Assessment + Timely Intervention + Parent Partnership + Student Accountability",
            Inches(0.6), Inches(4.5),
            Inches(12.13), Inches(1.5),
            font_size=Pt(14.5), color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 4  –  Characteristics of a Result-Oriented Teacher
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 4, "Teacher Profile")
slide_title(slide, "Characteristics of a Result-Oriented Teacher")

traits = [
    "Sets clear, measurable learning goals for every topic",
    "Plans lessons meticulously — no improvisation on delivery",
    "Tracks every student's progress week by week",
    "Uses data to make instructional decisions",
    "Takes complete ownership of student outcomes",
    "Provides timely, targeted remediation without delay",
]
# Two columns
mid = Inches(6.7)
for i, t in enumerate(traits[:3]):
    top_y = content_top() + i * Inches(1.25)
    add_rect(slide, Inches(0.55), top_y, Inches(0.45), Inches(0.45),
             fill_color=GOLD_DARK, line_color=None, line_width=Pt(0))
    add_textbox(slide, "✔", Inches(0.55), top_y, Inches(0.45), Inches(0.45),
                font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, t, Inches(1.1), top_y,
                mid - Inches(1.2), Inches(1.1),
                font_size=Pt(15), color=DARK_GRAY)

for i, t in enumerate(traits[3:]):
    top_y = content_top() + i * Inches(1.25)
    add_rect(slide, mid + Inches(0.1), top_y, Inches(0.45), Inches(0.45),
             fill_color=GOLD_DARK, line_color=None, line_width=Pt(0))
    add_textbox(slide, "✔", mid + Inches(0.1), top_y, Inches(0.45), Inches(0.45),
                font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, t, mid + Inches(0.65), top_y,
                Inches(6.0), Inches(1.1),
                font_size=Pt(15), color=DARK_GRAY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 5  –  The Power of Goal Setting
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 5, "Setting the North Star")
slide_title(slide, "The Power of Goal Setting")

# SMART box
highlight_box(slide, "S M A R T   G O A L S",
              Inches(0.55), content_top(),
              Inches(4.8), Inches(0.65),
              font_size=Pt(18), bold=True)

smart_items = [
    ("S", "Specific"),
    ("M", "Measurable"),
    ("A", "Achievable"),
    ("R", "Relevant"),
    ("T", "Time-bound"),
]
for i, (letter, word) in enumerate(smart_items):
    ty = content_top() + Inches(0.75) + i * Inches(0.72)
    add_rect(slide, Inches(0.55), ty, Inches(0.5), Inches(0.55),
             fill_color=GOLD_DARK, line_color=None, line_width=Pt(0))
    add_textbox(slide, letter, Inches(0.55), ty, Inches(0.5), Inches(0.55),
                font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, word, Inches(1.15), ty, Inches(3.7), Inches(0.55),
                font_size=Pt(15), color=DARK_GRAY)

# Goal types (right column)
add_textbox(slide, "Target Areas",
            Inches(6.0), content_top(),
            Inches(6.8), Inches(0.5),
            font_size=Pt(18), bold=True, color=NAVY)

goal_types = [
    ("📚", "Subject-wise targets per class"),
    ("🏫", "Class-wise average targets"),
    ("👤", "Student-wise individual targets"),
    ("🎯", "100% pass percentage"),
    ("📈", "Class average above 80%"),
    ("⚠", "Zero compartment / back cases"),
    ("↑",  "Minimum 10% improvement for every child"),
]
for i, (icon, text) in enumerate(goal_types):
    ty = content_top() + Inches(0.6) + i * Inches(0.62)
    add_textbox(slide, f"{icon}  {text}",
                Inches(6.0), ty, Inches(6.8), Inches(0.58),
                font_size=Pt(14), color=DARK_GRAY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 6  –  Reverse Planning Strategy
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 6, "Planning Backwards from Success")
slide_title(slide, "Reverse Planning Strategy")

stages = [
    ("Board Exam\n/ Final Goal", NAVY),
    ("Pre-Board\nExams", RGBColor(0x1F, 0x4E, 0x79)),
    ("Unit / Term\nTests", RGBColor(0x2E, 0x75, 0xB6)),
    ("Weekly\nAssessments", RGBColor(0x9D, 0xC3, 0xE6)),
    ("Daily Learning\nObjectives", RGBColor(0xD6, 0xE4, 0xF7)),
]
arrow_w = Inches(0.38)
box_w   = Inches(2.0)
total_w2 = len(stages)*box_w + (len(stages)-1)*arrow_w
start_l2 = (SLIDE_W - total_w2) / 2
box_h2   = Inches(2.0)
top_y2   = content_top() + Inches(0.4)

for i, (label, color) in enumerate(stages):
    lx = start_l2 + i*(box_w + arrow_w)
    add_rect(slide, lx, top_y2, box_w, box_h2,
             fill_color=color,
             line_color=GOLD_MID, line_width=Pt(1.8))
    add_textbox(slide, label, lx + Inches(0.05), top_y2 + Inches(0.55),
                box_w - Inches(0.1), Inches(1.1),
                font_size=Pt(13.5), bold=True,
                color=WHITE if i <= 2 else NAVY,
                align=PP_ALIGN.CENTER)
    if i < len(stages)-1:
        # Arrow
        ax = lx + box_w
        add_textbox(slide, "▶",
                    ax, top_y2 + Inches(0.7),
                    arrow_w, Inches(0.6),
                    font_size=Pt(20), color=GOLD_DARK,
                    align=PP_ALIGN.CENTER)

gold_divider(slide, top_y2 + box_h2 + Inches(0.25))

add_textbox(slide,
            "Plan backward from the Board Exam date.  Work inward to daily objectives so every lesson has purpose.",
            Inches(0.8), top_y2 + box_h2 + Inches(0.45),
            Inches(11.73), Inches(0.7),
            font_size=Pt(14.5), italic=True, color=NAVY,
            align=PP_ALIGN.CENTER)

add_textbox(slide,
            "Key Principle:  Each daily lesson is a brick in the wall of the Board Exam result.",
            Inches(0.8), top_y2 + box_h2 + Inches(1.15),
            Inches(11.73), Inches(0.5),
            font_size=Pt(13.5), color=DARK_GRAY,
            align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 7  –  Academic Monitoring Framework
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 7, "3-Tier Monitoring System")
slide_title(slide, "Academic Monitoring Framework")

cols = [
    ("🗓  Daily", NAVY, [
        "Learning objectives achieved?",
        "Homework completion rate",
        "Classroom participation",
        "Concept clarity checks",
    ]),
    ("📅  Weekly", RGBColor(0x14, 0x53, 0x78), [
        "Topic / chapter mastery",
        "Doubt resolution sessions",
        "Quick assessment results",
        "Weak student follow-up",
    ]),
    ("📆  Monthly", RGBColor(0x1F, 0x4E, 0x79), [
        "Performance analysis report",
        "Parent communication",
        "Remediation effectiveness",
        "Syllabus coverage audit",
    ]),
]

col_w = Inches(3.9)
col_h = Inches(4.15)
gap_c = Inches(0.27)
total_cw = 3*col_w + 2*gap_c
start_lc = (SLIDE_W - total_cw) / 2

for i, (heading, color, items) in enumerate(cols):
    lx = start_lc + i*(col_w + gap_c)
    add_rect(slide, lx, content_top(), col_w, Inches(0.6),
             fill_color=color, line_color=GOLD_MID, line_width=Pt(1.5))
    add_textbox(slide, heading, lx, content_top(),
                col_w, Inches(0.6),
                font_size=Pt(15.5), bold=True, color=GOLD_MID,
                align=PP_ALIGN.CENTER)
    add_rect(slide, lx, content_top() + Inches(0.6),
             col_w, col_h - Inches(0.6),
             fill_color=RGBColor(0xF5, 0xF0, 0xE0),
             line_color=GOLD_MID, line_width=Pt(1.5))
    for j, item in enumerate(items):
        ty = content_top() + Inches(0.7) + j * Inches(0.82)
        add_textbox(slide, f"▸  {item}",
                    lx + Inches(0.15), ty,
                    col_w - Inches(0.25), Inches(0.75),
                    font_size=Pt(13.5), color=DARK_GRAY)

highlight_box(slide,
              '"What gets measured gets managed — and what gets managed gets improved."',
              Inches(0.8), Inches(6.3),
              Inches(11.73), Inches(0.6),
              font_size=Pt(13.5))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 8  –  Data-Driven Teaching
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 8, "Using Evidence to Drive Instruction")
slide_title(slide, "Data-Driven Teaching")

areas = [
    ("📊", "Class Average", "Benchmark class performance against targets"),
    ("📖", "Chapter-wise\nPerformance", "Identify topics needing re-teaching"),
    ("👤", "Student-wise\nStrengths", "Celebrate wins, build confidence"),
    ("🔍", "Learning\nGaps", "Pinpoint exactly where students struggle"),
]

box_w3 = Inches(2.9)
box_h3 = Inches(2.6)
gap3   = Inches(0.27)
total3 = 4*box_w3 + 3*gap3
start3 = (SLIDE_W - total3) / 2

for i, (icon, heading, desc) in enumerate(areas):
    lx = start3 + i*(box_w3 + gap3)
    add_rect(slide, lx, content_top(), box_w3, box_h3,
             fill_color=RGBColor(0xFD, 0xF6, 0xE3),
             line_color=GOLD_DARK, line_width=Pt(2))
    add_textbox(slide, icon, lx, content_top() + Inches(0.15),
                box_w3, Inches(0.65),
                font_size=Pt(28), align=PP_ALIGN.CENTER)
    add_textbox(slide, heading, lx, content_top() + Inches(0.82),
                box_w3, Inches(0.75),
                font_size=Pt(14), bold=True, color=NAVY,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, desc, lx + Inches(0.1), content_top() + Inches(1.6),
                box_w3 - Inches(0.2), Inches(0.85),
                font_size=Pt(12.5), color=DARK_GRAY,
                align=PP_ALIGN.CENTER)

highlight_box(slide,
              '"What gets measured gets improved.  Data transforms effort into results."',
              Inches(1.2), Inches(4.55),
              Inches(10.93), Inches(0.7),
              font_size=Pt(14.5))

add_textbox(slide,
            "Action Step: Maintain a class data register — chapter-wise marks, attendance, homework, test scores.",
            Inches(0.8), Inches(5.45),
            Inches(11.73), Inches(0.6),
            font_size=Pt(13.5), italic=True, color=DARK_GRAY,
            align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 9  –  Identifying Learning Gaps
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 9, "Differentiated Support for Every Learner")
slide_title(slide, "Identifying Learning Gaps")

categories = [
    ("🏆", "High Achievers", NAVY, [
        "Enrichment activities",
        "Competitive exam preparation",
        "Leadership & peer-teaching roles",
        "Advanced problem sets",
    ]),
    ("📘", "Average Learners", RGBColor(0x1F, 0x4E, 0x79), [
        "Regular reinforcement exercises",
        "Concept revision cycles",
        "Structured practice sessions",
        "Motivational goal setting",
    ]),
    ("💡", "Slow Learners", RGBColor(0x2E, 0x75, 0xB6), [
        "Intensive one-on-one support",
        "Individual mentoring plan",
        "Simplified concept delivery",
        "Frequent positive feedback",
    ]),
]

cw = Inches(3.9); ch = Inches(4.0); cg = Inches(0.27)
tc = 3*cw + 2*cg; sc = (SLIDE_W - tc) / 2

for i, (icon, cat, color, items) in enumerate(categories):
    lx = sc + i*(cw+cg)
    add_rect(slide, lx, content_top(), cw, Inches(0.65),
             fill_color=color, line_color=GOLD_MID, line_width=Pt(1.5))
    add_textbox(slide, f"{icon}  {cat}",
                lx, content_top(), cw, Inches(0.65),
                font_size=Pt(14.5), bold=True, color=GOLD_MID,
                align=PP_ALIGN.CENTER)
    add_rect(slide, lx, content_top()+Inches(0.65),
             cw, ch-Inches(0.65),
             fill_color=RGBColor(0xF5, 0xF0, 0xE0),
             line_color=GOLD_MID, line_width=Pt(1.5))
    for j, item in enumerate(items):
        ty = content_top() + Inches(0.75) + j*Inches(0.78)
        add_textbox(slide, f"▸  {item}",
                    lx+Inches(0.15), ty, cw-Inches(0.25), Inches(0.72),
                    font_size=Pt(13), color=DARK_GRAY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 10  –  Remedial Action Plan
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 10, "Closing the Gap — Step by Step")
slide_title(slide, "Remedial Action Plan")

steps = [
    ("1", "IDENTIFY",  "📍", "Flag students scoring below threshold in each assessment."),
    ("2", "DIAGNOSE",  "🔎", "Analyse root cause — concept gap, practice deficit, or motivation issue."),
    ("3", "INTERVENE", "🛠", "Design targeted remedial sessions, worksheets, and peer tutoring."),
    ("4", "MONITOR",   "👁", "Track progress weekly; confirm improvement or escalate support."),
    ("5", "REVIEW",    "✅", "Validate effectiveness; celebrate progress and re-set targets."),
]

bw = Inches(2.3); bh = Inches(2.75); bg_ = Inches(0.18)
total_bw = len(steps)*bw + (len(steps)-1)*bg_
start_bl = (SLIDE_W - total_bw) / 2

for i, (num, label, icon, desc) in enumerate(steps):
    lx = start_bl + i*(bw+bg_)
    col_bg = [NAVY, RGBColor(0x1F,0x4E,0x79), RGBColor(0x2E,0x75,0xB6),
              RGBColor(0x4A,0x86,0xC8), RGBColor(0x9D,0xC3,0xE6)][i]
    add_rect(slide, lx, content_top(), bw, Inches(0.55),
             fill_color=col_bg, line_color=GOLD_MID, line_width=Pt(1.5))
    add_textbox(slide, f"Step {num}", lx, content_top(),
                bw, Inches(0.55), font_size=Pt(12), bold=True,
                color=GOLD_LIGHT, align=PP_ALIGN.CENTER)
    add_rect(slide, lx, content_top()+Inches(0.55),
             bw, bh-Inches(0.55),
             fill_color=RGBColor(0xFD,0xF6,0xE3),
             line_color=GOLD_MID, line_width=Pt(1.5))
    add_textbox(slide, icon, lx, content_top()+Inches(0.65),
                bw, Inches(0.7), font_size=Pt(24), align=PP_ALIGN.CENTER)
    add_textbox(slide, label, lx, content_top()+Inches(1.35),
                bw, Inches(0.55), font_size=Pt(14), bold=True,
                color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, desc, lx+Inches(0.1), content_top()+Inches(1.9),
                bw-Inches(0.2), Inches(0.85), font_size=Pt(11.5),
                color=DARK_GRAY, align=PP_ALIGN.CENTER)

highlight_box(slide,
              "Remediation must be  IMMEDIATE  ·  SPECIFIC  ·  MEASURABLE",
              Inches(1.5), Inches(5.25), Inches(10.33), Inches(0.65),
              font_size=Pt(14.5), bold=True)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 11  –  Classroom Strategies
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 11, "Evidence-Based Instructional Strategies")
slide_title(slide, "Classroom Strategies for Better Results")

strategies = [
    ("🔄", "Active Learning",    "Students engage, discuss, and apply — not just listen."),
    ("👥", "Peer Teaching",      "Students who teach retain 90% of the material."),
    ("🧠", "Retrieval Practice", "Regular low-stakes quizzes dramatically boost retention."),
    ("🗺", "Concept Mapping",    "Visual connections deepen conceptual understanding."),
    ("✏", "Regular Quizzes",    "Frequent testing reveals gaps before they widen."),
    ("⬆", "HOTS Questions",     "Higher-order thinking prepares for board-level questions."),
]

sw = Inches(3.9); sh = Inches(1.55); sg = Inches(0.22)
rows = [(0,1,2), (3,4,5)]
for row_i, row in enumerate(rows):
    for col_i, idx in enumerate(row):
        icon, strat, desc = strategies[idx]
        lx = Inches(0.55) + col_i*(sw+sg)
        ty = content_top() + row_i*(sh+sg)
        add_rect(slide, lx, ty, sw, sh,
                 fill_color=RGBColor(0xFD,0xF6,0xE3),
                 line_color=GOLD_DARK, line_width=Pt(1.8))
        add_textbox(slide, f"{icon}  {strat}",
                    lx+Inches(0.1), ty+Inches(0.1),
                    sw-Inches(0.2), Inches(0.55),
                    font_size=Pt(14.5), bold=True, color=NAVY)
        add_textbox(slide, desc,
                    lx+Inches(0.1), ty+Inches(0.65),
                    sw-Inches(0.2), Inches(0.75),
                    font_size=Pt(12.5), color=DARK_GRAY)

highlight_box(slide,
              "Combine strategies: Teach → Quiz → Reteach → Map → Discuss → Apply",
              Inches(1.0), Inches(5.5), Inches(11.33), Inches(0.65),
              font_size=Pt(14.5))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 12  –  Question Paper Analysis
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 12, "Decode the Exam — Teach with Precision")
slide_title(slide, "Question Paper Analysis")

left_items = [
    ("📊", "Analyse Question Trends", "Identify which topics appear year after year."),
    ("📌", "Frequently Asked Concepts", "Ensure these are taught thoroughly with examples."),
    ("🧩", "Competency-Based Questions", "Align teaching to application and analysis levels."),
    ("🔺", "Bloom's Taxonomy Distribution", "Balance Remember, Understand, Apply, Analyse, Evaluate, Create."),
]

for i, (icon, heading, desc) in enumerate(left_items):
    ty = content_top() + i*Inches(1.15)
    add_rect(slide, Inches(0.55), ty+Inches(0.05),
             Inches(0.5), Inches(0.5),
             fill_color=GOLD_DARK, line_color=None, line_width=Pt(0))
    add_textbox(slide, icon,
                Inches(0.55), ty+Inches(0.05),
                Inches(0.5), Inches(0.5),
                font_size=Pt(18), align=PP_ALIGN.CENTER)
    add_textbox(slide, heading,
                Inches(1.2), ty, Inches(11.0), Inches(0.5),
                font_size=Pt(15.5), bold=True, color=NAVY)
    add_textbox(slide, desc,
                Inches(1.2), ty+Inches(0.5), Inches(11.0), Inches(0.55),
                font_size=Pt(13.5), color=DARK_GRAY)

highlight_box(slide,
              '"Teach what is tested.  Test what was taught.  Align every lesson to the exam blueprint."',
              Inches(0.8), Inches(6.1), Inches(11.73), Inches(0.7),
              font_size=Pt(13.5))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 13  –  Assessment for Learning
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 13, "Assessment AS Learning — Not Just OF Learning")
slide_title(slide, "Assessment for Learning")

# Left: purpose
add_textbox(slide, "Assessment Should…",
            Inches(0.55), content_top(), Inches(6.0), Inches(0.5),
            font_size=Pt(17), bold=True, color=NAVY)
purposes = [
    "Guide and adjust teaching in real time",
    "Identify individual weaknesses early",
    "Improve student understanding progressively",
    "Inform the next lesson plan",
]
add_bullet_block(slide, purposes,
                 Inches(0.55), content_top()+Inches(0.55),
                 Inches(5.8), Inches(3.0),
                 font_size=Pt(15), bullet_char="✦", color=DARK_GRAY)

# Divider
add_rect(slide, Inches(6.6), content_top(), Pt(2), Inches(4.5),
         fill_color=GOLD_MID, line_color=None, line_width=Pt(0))

# Right: tools
add_textbox(slide, "Assessment Tools",
            Inches(6.9), content_top(), Inches(6.0), Inches(0.5),
            font_size=Pt(17), bold=True, color=NAVY)

tools = [
    ("📝", "Exit Tickets",       "Quick end-of-lesson insight"),
    ("⚡", "Quick Checks",       "Mid-lesson comprehension pulse"),
    ("📋", "Weekly Tests",       "Chapter mastery verification"),
    ("🗣", "Oral Assessments",   "Speaking & reasoning skills"),
]
for i, (icon, tool, desc) in enumerate(tools):
    ty = content_top() + Inches(0.6) + i*Inches(0.92)
    add_rect(slide, Inches(6.9), ty, Inches(5.8), Inches(0.8),
             fill_color=RGBColor(0xFD,0xF6,0xE3),
             line_color=GOLD_MID, line_width=Pt(1.2))
    add_textbox(slide, f"{icon}  {tool}  —  {desc}",
                Inches(7.0), ty+Inches(0.1),
                Inches(5.6), Inches(0.65),
                font_size=Pt(13.5), color=DARK_GRAY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 14  –  Parent Partnership
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 14, "Engaging the Third Teacher")
slide_title(slide, "Parent Partnership")

# Central visual
add_rect(slide, Inches(5.4), content_top()+Inches(0.1),
         Inches(2.5), Inches(2.5),
         fill_color=NAVY, line_color=GOLD_MID, line_width=Pt(2.5))
add_textbox(slide, "🤝\nPartnership\nZone",
            Inches(5.4), content_top()+Inches(0.25),
            Inches(2.5), Inches(2.2),
            font_size=Pt(14), bold=True, color=GOLD_MID,
            align=PP_ALIGN.CENTER)

parent_know = [
    "Child's academic strengths",
    "Areas requiring improvement",
    "Attendance & punctuality record",
    "Personalised study plan",
]
add_textbox(slide, "Parents Should Know…",
            Inches(0.4), content_top(), Inches(4.7), Inches(0.5),
            font_size=Pt(16), bold=True, color=NAVY)
add_bullet_block(slide, parent_know,
                 Inches(0.4), content_top()+Inches(0.55),
                 Inches(4.7), Inches(2.5),
                 font_size=Pt(14.5), bullet_char="▸", color=DARK_GRAY)

teacher_do = [
    "Monthly progress reports",
    "PTM with actionable feedback",
    "WhatsApp / diary communication",
    "Home-study guidance",
]
add_textbox(slide, "Teachers Should Do…",
            Inches(8.2), content_top(), Inches(4.7), Inches(0.5),
            font_size=Pt(16), bold=True, color=NAVY)
add_bullet_block(slide, teacher_do,
                 Inches(8.2), content_top()+Inches(0.55),
                 Inches(4.7), Inches(2.5),
                 font_size=Pt(14.5), bullet_char="▸", color=DARK_GRAY)

highlight_box(slide,
              "Strong Parent + Teacher Collaboration  =  Dramatically Better Results",
              Inches(1.0), Inches(5.1), Inches(11.33), Inches(0.75),
              font_size=Pt(16), bold=True)

add_textbox(slide,
            "Involve parents as partners, not just observers.  Their support at home multiplies classroom effort.",
            Inches(0.8), Inches(6.05), Inches(11.73), Inches(0.55),
            font_size=Pt(13), italic=True, color=DARK_GRAY,
            align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 15  –  Building Student Ownership
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 15, "Empowering Learners")
slide_title(slide, "Building Student Ownership")

ownership = [
    ("🎯", "Set Personal Goals",
     "Students define their own subject targets each month."),
    ("📈", "Track Their Progress",
     "Maintain a personal performance tracker / chart."),
    ("🔍", "Analyse Mistakes",
     "Review every wrong answer — understand the error, not just the correction."),
    ("💭", "Reflect Regularly",
     "Weekly self-reflection: What did I master? Where do I need more effort?"),
]

bw4 = Inches(5.85); bh4 = Inches(1.6); bg4 = Inches(0.22)
for i, (icon, heading, desc) in enumerate(ownership):
    row = i // 2; col = i % 2
    lx = Inches(0.45) + col*(bw4+bg4)
    ty = content_top() + row*(bh4+bg4)
    add_rect(slide, lx, ty, bw4, bh4,
             fill_color=RGBColor(0xFD,0xF6,0xE3),
             line_color=GOLD_DARK, line_width=Pt(1.8))
    add_textbox(slide, icon, lx+Inches(0.15), ty+Inches(0.12),
                Inches(0.6), Inches(0.6), font_size=Pt(26))
    add_textbox(slide, heading, lx+Inches(0.85), ty+Inches(0.1),
                bw4-Inches(1.0), Inches(0.55),
                font_size=Pt(15), bold=True, color=NAVY)
    add_textbox(slide, desc, lx+Inches(0.85), ty+Inches(0.65),
                bw4-Inches(1.0), Inches(0.8),
                font_size=Pt(13.5), color=DARK_GRAY)

highlight_box(slide,
              '"Teach students to own their learning — a self-driven student outperforms a pushed student every time."',
              Inches(0.8), Inches(5.5), Inches(11.73), Inches(0.75),
              font_size=Pt(13.5))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 16  –  5-STAR Model
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 16, "The Result Improvement Cycle")
slide_title(slide, "The 5-STAR Result Improvement Model")

stars = [
    ("S", "Set Goals",         "Define clear, measurable targets per student, per subject."),
    ("T", "Track Progress",    "Weekly data collection and trend analysis."),
    ("A", "Analyse Data",      "Deep-dive into performance patterns and gaps."),
    ("R", "Remediate",         "Targeted, immediate intervention for every gap."),
    ("S", "Sustain Excellence","Review, celebrate, and raise the bar continuously."),
]

sw5 = Inches(2.3); sh5 = Inches(3.2); sg5 = Inches(0.2)
tw5 = 5*sw5 + 4*sg5; sl5 = (SLIDE_W-tw5)/2

colors5 = [NAVY, RGBColor(0x1F,0x4E,0x79), RGBColor(0x2E,0x75,0xB6),
           RGBColor(0x4A,0x86,0xC8), RGBColor(0x9D,0xC3,0xE6)]

for i, (letter, label, desc) in enumerate(stars):
    lx = sl5 + i*(sw5+sg5)
    add_rect(slide, lx, content_top(), sw5, sh5,
             fill_color=colors5[i],
             line_color=GOLD_MID, line_width=Pt(2))
    # Gold star letter
    add_rect(slide, lx + (sw5-Inches(0.75))/2,
             content_top()+Inches(0.2),
             Inches(0.75), Inches(0.75),
             fill_color=GOLD_DARK, line_color=None, line_width=Pt(0))
    add_textbox(slide, letter,
                lx + (sw5-Inches(0.75))/2,
                content_top()+Inches(0.2),
                Inches(0.75), Inches(0.75),
                font_size=Pt(24), bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, label, lx+Inches(0.1),
                content_top()+Inches(1.1), sw5-Inches(0.2), Inches(0.75),
                font_size=Pt(13.5), bold=True, color=GOLD_LIGHT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, desc, lx+Inches(0.1),
                content_top()+Inches(1.9), sw5-Inches(0.2), Inches(1.1),
                font_size=Pt(11.5), color=WHITE,
                align=PP_ALIGN.CENTER)

add_textbox(slide, "⭐  S  –  T  –  A  –  R  –  S  ⭐",
            Inches(2.5), content_top()+sh5+Inches(0.2),
            Inches(8.33), Inches(0.5),
            font_size=Pt(20), bold=True, color=GOLD_DARK,
            align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 17  –  Teacher Accountability Matrix
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 17, "Professional Accountability")
slide_title(slide, "Teacher Accountability Matrix")

headers = ["Metric", "Frequency", "Owner", "Target"]
rows17  = [
    ["Student Attendance",      "Daily",    "Class Teacher", "≥ 95%"],
    ["Homework Completion",     "Daily",    "Subject Teacher","100%"],
    ["Test / Assessment Score", "Weekly",   "Subject Teacher","Class Avg ≥ 75%"],
    ["Weak Student Follow-up",  "Weekly",   "Subject Teacher","Zero Left Behind"],
    ["Parent Interactions",     "Monthly",  "Class Teacher", "Every Parent"],
    ["Improvement Percentage",  "Monthly",  "Subject Teacher","≥ 10% per student"],
]

col_widths = [Inches(3.8), Inches(2.2), Inches(3.0), Inches(3.0)]
col_starts = [Inches(0.35)]
for cw7 in col_widths[:-1]:
    col_starts.append(col_starts[-1] + cw7 + Inches(0.04))

row_h = Inches(0.52); head_h = Inches(0.55)
top17 = content_top()

# Header row
for ci, (hdr, lx, cw7) in enumerate(zip(headers, col_starts, col_widths)):
    add_rect(slide, lx, top17, cw7, head_h,
             fill_color=NAVY, line_color=GOLD_MID, line_width=Pt(1.5))
    add_textbox(slide, hdr, lx, top17, cw7, head_h,
                font_size=Pt(14), bold=True, color=GOLD_MID,
                align=PP_ALIGN.CENTER)

# Data rows
for ri, row in enumerate(rows17):
    ty = top17 + head_h + ri*row_h
    bg = RGBColor(0xFD,0xF6,0xE3) if ri%2==0 else RGBColor(0xF0,0xE8,0xC8)
    for ci, (cell, lx, cw7) in enumerate(zip(row, col_starts, col_widths)):
        add_rect(slide, lx, ty, cw7, row_h,
                 fill_color=bg, line_color=GOLD_LIGHT, line_width=Pt(0.8))
        add_textbox(slide, cell, lx+Inches(0.1), ty+Inches(0.05),
                    cw7-Inches(0.15), row_h-Inches(0.05),
                    font_size=Pt(13), color=DARK_GRAY,
                    align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 18  –  Monthly Review Mechanism
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 18, "Continuous Improvement Cycle")
slide_title(slide, "Monthly Review Mechanism")

reviews = [
    ("📊", "Subject Performance",      "Analyse subject-wise class average vs. target."),
    ("🏫", "Class Performance",        "Compare class progress with school benchmarks."),
    ("📈", "Improvement Trends",       "Track month-on-month growth for each student."),
    ("🛠", "Intervention Effectiveness","Evaluate whether remedial actions are working."),
]

rw = Inches(5.85); rh = Inches(1.65); rg = Inches(0.22)
for i, (icon, heading, desc) in enumerate(reviews):
    row = i//2; col = i%2
    lx = Inches(0.45) + col*(rw+rg)
    ty = content_top() + row*(rh+rg)
    add_rect(slide, lx, ty, rw, rh,
             fill_color=RGBColor(0xFD,0xF6,0xE3),
             line_color=GOLD_DARK, line_width=Pt(1.8))
    add_textbox(slide, icon, lx+Inches(0.15), ty+Inches(0.12),
                Inches(0.7), Inches(0.7), font_size=Pt(28))
    add_textbox(slide, heading, lx+Inches(0.95), ty+Inches(0.12),
                rw-Inches(1.05), Inches(0.55),
                font_size=Pt(15), bold=True, color=NAVY)
    add_textbox(slide, desc, lx+Inches(0.95), ty+Inches(0.7),
                rw-Inches(1.05), Inches(0.8),
                font_size=Pt(13.5), color=DARK_GRAY)

add_textbox(slide, "Review Cycle:  Analyse  →  Decide  →  Act  →  Monitor  →  Repeat",
            Inches(0.8), Inches(5.45), Inches(11.73), Inches(0.55),
            font_size=Pt(14.5), bold=True, color=NAVY,
            align=PP_ALIGN.CENTER)
highlight_box(slide,
              "Schedule a fixed Monthly Review Day — make it a non-negotiable calendar event.",
              Inches(1.2), Inches(6.1), Inches(10.93), Inches(0.65),
              font_size=Pt(13.5))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 19  –  High-Performance Habits
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 19, "What Top Teachers Do Differently")
slide_title(slide, "High-Performance Habits of Top Teachers")

habits = [
    ("☀", "Preparation Before Class",
     "Lesson plan is finalised, resources ready, questions pre-identified."),
    ("📝", "Feedback After Class",
     "Reflect on what worked, what didn't, and what to adjust tomorrow."),
    ("🔁", "Consistency",
     "Show up with the same energy and standards every single day."),
    ("🪞", "Reflection",
     "Weekly review: Am I achieving my targets? Where are my students stuck?"),
    ("📚", "Continuous Learning",
     "Stay updated with pedagogy, subject developments, and exam patterns."),
]

hw = Inches(11.7); hh = Inches(0.9); hg = Inches(0.18)
start_h = content_top()

for i, (icon, habit, desc) in enumerate(habits):
    ty = start_h + i*(hh+hg)
    add_rect(slide, Inches(0.45), ty, Inches(0.7), hh,
             fill_color=GOLD_DARK, line_color=None, line_width=Pt(0))
    add_textbox(slide, icon, Inches(0.45), ty,
                Inches(0.7), hh, font_size=Pt(22), align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(1.2), ty, hw-Inches(0.75), hh,
             fill_color=RGBColor(0xFD,0xF6,0xE3) if i%2==0 else RGBColor(0xF0,0xE8,0xC8),
             line_color=GOLD_LIGHT, line_width=Pt(0.8))
    add_textbox(slide, f"{habit}  —  {desc}",
                Inches(1.35), ty+Inches(0.15),
                hw-Inches(0.95), hh-Inches(0.2),
                font_size=Pt(14), color=DARK_GRAY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 20  –  School Excellence Framework
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 20, "The Ecosystem of Excellence")
slide_title(slide, "School Excellence Framework")

# Centre circle
cx = SLIDE_W/2; cy = Inches(4.0)
add_rect(slide, cx-Inches(1.25), cy-Inches(0.75),
         Inches(2.5), Inches(1.5),
         fill_color=GOLD_DARK, line_color=GOLD_MID, line_width=Pt(2))
add_textbox(slide, "EXCELLENT\nRESULTS",
            cx-Inches(1.25), cy-Inches(0.65),
            Inches(2.5), Inches(1.2),
            font_size=Pt(14), bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)

pillars = [
    ("🏛", "Visionary\nLeadership",  Inches(1.4),  Inches(2.8)),
    ("🎓", "Effective\nTeachers",    Inches(4.15), Inches(2.8)),
    ("🙋", "Engaged\nStudents",      Inches(7.15), Inches(2.8)),
    ("👨‍👩‍👧", "Supportive\nParents",  Inches(9.9),  Inches(2.8)),
    ("📊", "Continuous\nMonitoring", Inches(5.65), Inches(5.5)),
]

for icon, label, lx, ty in pillars:
    add_rect(slide, lx, ty, Inches(2.3), Inches(1.55),
             fill_color=NAVY, line_color=GOLD_MID, line_width=Pt(1.8))
    add_textbox(slide, f"{icon}\n{label}", lx, ty,
                Inches(2.3), Inches(1.55),
                font_size=Pt(13), bold=True, color=GOLD_LIGHT,
                align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 21  –  Action Plan
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
decorate(slide, 21, "Your 90-Day Execution Roadmap")
slide_title(slide, "Action Plan — Starting Today")

tasks = [
    ("1", "Define Subject Target",        "Set class average and pass-% targets for the term."),
    ("2", "Set Class Target",             "Benchmark current performance; plan to exceed it."),
    ("3", "Prepare Weak Student List",    "Identify and categorise every at-risk student."),
    ("4", "Design Remedial Plan",         "Map specific interventions for each category."),
    ("5", "Create Monitoring Schedule",   "Fix weekly and monthly review dates in the calendar."),
]

for i, (num, task, desc) in enumerate(tasks):
    ty = content_top() + i*Inches(0.95)
    # Number badge
    add_rect(slide, Inches(0.45), ty+Inches(0.08),
             Inches(0.55), Inches(0.55),
             fill_color=GOLD_DARK, line_color=None, line_width=Pt(0))
    add_textbox(slide, num, Inches(0.45), ty+Inches(0.08),
                Inches(0.55), Inches(0.55),
                font_size=Pt(16), bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    # Task bar
    add_rect(slide, Inches(1.1), ty,
             Inches(11.7), Inches(0.85),
             fill_color=RGBColor(0xFD,0xF6,0xE3) if i%2==0 else RGBColor(0xF0,0xE8,0xC8),
             line_color=GOLD_LIGHT, line_width=Pt(0.8))
    add_textbox(slide, f"{task}",
                Inches(1.25), ty+Inches(0.04),
                Inches(4.5), Inches(0.45),
                font_size=Pt(14.5), bold=True, color=NAVY)
    add_textbox(slide, desc,
                Inches(5.85), ty+Inches(0.04),
                Inches(6.8), Inches(0.78),
                font_size=Pt(13.5), color=DARK_GRAY)

highlight_box(slide,
              "Each teacher leaves today with a clear, written Action Plan — not just inspiration.",
              Inches(0.8), Inches(6.2), Inches(11.73), Inches(0.65),
              font_size=Pt(14))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 22  –  Commitment / Closing
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, NAVY)

# Borders
m = Inches(0.12)
add_rect(slide, m, m, SLIDE_W-2*m, SLIDE_H-2*m,
         fill_color=None, line_color=GOLD_MID, line_width=Pt(4))
add_rect(slide, Inches(0.28), Inches(0.28),
         SLIDE_W-Inches(0.56), SLIDE_H-Inches(0.56),
         fill_color=None, line_color=GOLD_LIGHT, line_width=Pt(1.2))
corner_ornament(slide)

gold_divider(slide, Inches(1.3), 0.7)
gold_divider(slide, Inches(6.15), 0.7)

add_textbox(slide, "✦  Teacher's Pledge  ✦",
            Inches(1), Inches(0.5), Inches(11.33), Inches(0.65),
            font_size=Pt(20), bold=True, color=GOLD_MID,
            align=PP_ALIGN.CENTER)

pledge = (
    '"I will teach with PURPOSE,\n'
    'assess with INTEGRITY,\n'
    'monitor with DILIGENCE,\n'
    'and support every learner\n'
    'to achieve their highest potential."'
)
add_textbox(slide, pledge,
            Inches(1.2), Inches(1.5), Inches(10.93), Inches(3.0),
            font_name="Palatino Linotype",
            font_size=Pt(21), italic=True, color=WHITE,
            align=PP_ALIGN.CENTER)

add_textbox(slide,
            '"Every Child Matters  ·  Every Mark Matters  ·  Every Day Matters"',
            Inches(1), Inches(4.6), Inches(11.33), Inches(0.65),
            font_size=Pt(16.5), bold=True, color=GOLD_MID,
            align=PP_ALIGN.CENTER)

add_textbox(slide,
            '"Great results are not achieved by chance;\n'
            'they are achieved by design, discipline, and dedication."',
            Inches(1.5), Inches(5.35), Inches(10.33), Inches(1.0),
            font_size=Pt(14), italic=True, color=GOLD_LIGHT,
            align=PP_ALIGN.CENTER)

add_textbox(slide,
            "From Effort to Excellence  •  Teacher Training Programme",
            Inches(1), Inches(6.9), Inches(11.33), Inches(0.4),
            font_size=Pt(10), color=GOLD_LIGHT,
            italic=True, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out_path = "/home/user/bng_website/Teacher_Training_Excellence_PPT.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
